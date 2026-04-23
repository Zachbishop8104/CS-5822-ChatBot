import os
import re
import io
from pathlib import Path
from fastapi import UploadFile

BASE_DIR = Path(__file__).resolve().parent.parent

def extract_text_from_bytes(filename: str, contents: bytes) -> str:
    suffix = Path(filename).suffix.lower()

    if suffix == ".txt":
        return contents.decode("utf-8", errors="ignore")

    elif suffix == ".pdf":
        import fitz
        doc = fitz.open(stream=contents, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    elif suffix == ".docx":
        from docx import Document
        doc = Document(io.BytesIO(contents))
        return "\n".join(p.text for p in doc.paragraphs)

    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(contents))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape_text = shape.text.strip()
                    # Skip shapes that look like diagram notation
                    # (short text with no real words)
                    words = shape_text.split()
                    real_words = sum(1 for w in words if re.match(r'[a-zA-Z]{3,}', w))
                    if len(words) > 0 and real_words / len(words) < 0.4:
                        continue
                    lines.append(shape_text)
        return "\n".join(lines)

    elif suffix == ".rtf":
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(contents.decode("utf-8", errors="ignore"))

    elif suffix == ".md":
        text = contents.decode("utf-8", errors="ignore")
        text = re.sub(r"[#*`_~>\[\]]+", "", text)
        return text

    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def normalize_text(text: str) -> str:
    # Fix hyphenated line breaks
    text = re.sub(r"-\n\s*", "", text)
    # Join lines that are part of the same sentence
    text = re.sub(r"(?<=[a-z])\n(?=[a-z])", " ", text)
    # Strip bullet characters at line starts
    text = re.sub(r"(?m)^[\s]*[•\-\*]\s*", "", text)
    # Collapse multiple spaces
    text = re.sub(r"[ \t]+", " ", text)
    # Collapse excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


async def download_note(file: UploadFile, user_id: str):
    contents = await file.read()

    user_dir = BASE_DIR / "users" / user_id
    os.makedirs(user_dir, exist_ok=True)

    # save only the normalized .txt
    txt_filename = Path(file.filename).stem + ".txt"
    txt_path = user_dir / txt_filename

    if os.path.exists(txt_path):
        return {"error": "File already exists"}

    try:
        text = normalize_text(extract_text_from_bytes(file.filename, contents))
    except ValueError as e:
        return {"error": str(e)}

    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

    return "File uploaded successfully"